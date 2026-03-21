from __future__ import annotations

from pathlib import Path

from homie_core.rag.parsers import ParsedDocument, TextBlock, register_parser


@register_parser("pptx")
def parse_pptx(path: Path) -> ParsedDocument:
    try:
        from pptx import Presentation
    except ImportError:
        return ParsedDocument(source_path=str(path), metadata={"error": "python-pptx not installed"})
    prs = Presentation(str(path))
    blocks = []
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text.strip())
        if texts:
            blocks.append(TextBlock(content="\n".join(texts), block_type="paragraph", page=slide_num))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                blocks.append(
                    TextBlock(content=f"[Speaker Notes] {notes}", block_type="paragraph", page=slide_num)
                )
    return ParsedDocument(
        text_blocks=blocks,
        metadata={"format": "pptx", "slides": len(prs.slides)},
        source_path=str(path),
    )
